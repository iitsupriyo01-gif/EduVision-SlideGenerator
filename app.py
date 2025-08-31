import streamlit as st
import os
from openai import OpenAI
from pptx import Presentation
import io

# Load API key
api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)

st.title("EduVision – Automated Slide Generator")
st.write("Upload your notes (TXT/Markdown) and generate slides automatically.")

uploaded_file = st.file_uploader("Upload your notes", type=["txt", "md"])

def create_pptx_from_outline(outline_text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout

    slides = outline_text.split("Slide ")
    for slide in slides:
        if not slide.strip():
            continue
        lines = slide.strip().split("\n")
        if len(lines) < 2:
            continue
        title = lines[0].replace(":", "").strip()
        bullets = [line.strip("- ").strip() for line in lines[1:] if line.strip()]

        slide_obj = prs.slides.add_slide(slide_layout)
        slide_obj.shapes.title.text = title
        content = slide_obj.placeholders[1].text_frame
        content.clear()  # Remove placeholder text

        for bullet in bullets:
            p = content.add_paragraph()
            p.text = bullet
            p.level = 0  # Makes it a bullet point

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

if uploaded_file is not None:
    content = uploaded_file.read().decode("utf-8")
    st.text_area("Preview of Notes", content, height=200)

    if st.button("Generate Slide Outline"):
        with st.spinner("Generating slides... please wait..."):
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "You are an assistant that creates slide decks."},
                        {"role": "user", "content": f"Create 5 slides from these notes. Give a title and 3-4 bullet points per slide.\n\n{content}"}
                    ]
                )
                outline = response.choices[0].message.content
                st.subheader("Generated Slide Outline")
                st.write(outline)

                # Generate PPTX and add download button
                pptx_file = create_pptx_from_outline(outline)
                st.download_button(
                    label="Download PowerPoint",
                    data=pptx_file,
                    file_name="EduVision_Slides.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"Error: {e}")
